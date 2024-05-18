from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "OCA"
subtitle.text = "welcome"

def count_text_in_pptx(filename):
    try:
        prs = Presentation(filename)
        total_text_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        total_text_count += len(paragraph.text.split())
        return total_text_count
    except Exception as e:
        print("エラーが発生しました:", e)

# PowerPointファイル名を指定してテキスト数をカウント
filename = "test.pptx"  # 実際のファイル名に置き換えてください
text_count = count_text_in_pptx(filename)
if text_count is not None:
    print(f"{filename} 内のテキスト数: {text_count}")

prs.save('test2.pptx')